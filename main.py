from fastapi import FastAPI, UploadFile, Form, Response
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
import io
import json
import logging
import time

app = FastAPI()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
)
logger = logging.getLogger("pptx-generator-api")

# 1. CORS Configuration: Allow your SPFx web part to communicate with this server
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # For production, replace "*" with "https://yourtenant.sharepoint.com"
    allow_methods=["POST"],
    allow_headers=["*"],
)

def delete_slide(prs, index):
    """Helper function to remove a slide by its index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])

def normalize_pptx_filename(value):
    """Return a safe .pptx filename suitable for Content-Disposition."""
    fallback_name = "generated_deck.pptx"

    if not isinstance(value, str):
        return fallback_name

    candidate = value.strip()
    if not candidate:
        return fallback_name

    # Keep only the final segment if a path-like value is sent.
    candidate = candidate.split("/")[-1].split("\\")[-1]

    sanitized = "".join(
        "_" if (char in '<>:"/\\|?*' or ord(char) < 32) else char
        for char in candidate
    ).rstrip(" .")

    if not sanitized:
        return fallback_name

    if not sanitized.lower().endswith(".pptx"):
        sanitized = f"{sanitized}.pptx"

    return sanitized

@app.post("/generate-document")
async def generate_document(
    file: UploadFile, 
    data: str = Form(...) # The JSON string sent from SPFx containing variables and instructions
):
    started_at = time.perf_counter()
    logger.info("Request received at /generate-document")
    logger.info("Incoming file name: %s", file.filename)

    # 2. Parse the incoming JSON data
    instructions = json.loads(data)
    variables = instructions.get("variables", {})
    replacement_variables = dict(variables) if isinstance(variables, dict) else {}
    # fileName is control metadata and should not be used as a text replacement key.
    raw_file_name = replacement_variables.pop("fileName", None)
    if raw_file_name is None:
        raw_file_name = replacement_variables.get("{{FILE_NAME}}")
    response_file_name = normalize_pptx_filename(raw_file_name)
    slides_to_remove = instructions.get("deleteSlides", [])
    metadata = instructions.get("metadata", {})

    logger.info("Variables count: %s", len(variables))
    logger.info("Slides requested to delete: %s", slides_to_remove)
    logger.info("Response file name: %s", response_file_name)
    logger.info("Metadata: %s", metadata)

    # 3. Read the incoming file (Blob) into server memory
    file_bytes = await file.read()
    logger.info("Input file size (bytes): %s", len(file_bytes))
    input_stream = io.BytesIO(file_bytes)
    prs = Presentation(input_stream)
    logger.info("Loaded presentation with %s slides", len(prs.slides))

    # 4. Delete specified slides (Must be done in reverse order to avoid index shifting)
    deleted_count = 0
    for index in sorted(slides_to_remove, reverse=True):
        if index < len(prs.slides):
            delete_slide(prs, index)
            deleted_count += 1
            logger.info("Deleted slide at index: %s", index)
        else:
            logger.warning("Skipped invalid slide index: %s", index)
    logger.info("Total deleted slides: %s", deleted_count)

    # 5. Search and Replace text in shapes and tables
    replacement_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            # Handle standard text boxes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacement_variables.items():
                            if key in run.text:
                                run.text = run.text.replace(key, str(value))
                                replacement_count += 1
            
            # Handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for key, value in replacement_variables.items():
                                    if key in run.text:
                                        run.text = run.text.replace(key, str(value))
                                        replacement_count += 1
    logger.info("Total text replacements applied: %s", replacement_count)

    # 6. Save the modified presentation to a new memory buffer
    output_stream = io.BytesIO()
    prs.save(output_stream)
    logger.info("Presentation saved to memory buffer")
    
    # Reset the buffer's "cursor" to the beginning before reading it out
    output_stream.seek(0)
    output_bytes = output_stream.getvalue()

    elapsed_ms = int((time.perf_counter() - started_at) * 1000)
    logger.info("Returning modified presentation. Output size: %s bytes", len(output_bytes))
    logger.info("Request processing completed in %s ms", elapsed_ms)

    # 7. Return the raw binary stream (Blob) directly back to SPFx
    return Response(
        content=output_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f'attachment; filename="{response_file_name}"'
        },
    )

@app.get("/")
async def root():
    return {"message": "Welcome to the PowerPoint Generator API"}

@app.get("/owner")
async def get_owner():
    return {"owner": "nijobair"}